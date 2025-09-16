from celery import Celery
from transformers import pipeline
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation

celery_app = Celery('tasks', broker='redis://localhost:6379/0')

# Hugging Face summarization pipeline
summarizer = pipeline("summarization", model="t5-small")

@celery_app.task
def process_file(file_path, file_type):
    # Extract text and run summarization based on file type
    if file_type == 'pdf':
        pdf_reader = PyPDF2.PdfReader(file_path)
        text = "".join([page.extract_text() for page in pdf_reader.pages])
    
    elif file_type == 'docx':
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    
    elif file_type == 'pptx':
        ppt = Presentation(file_path)
        text = "\n".join([slide.shapes.title.text for slide in ppt.slides if slide.shapes.title])
    
    elif file_type in ['csv', 'xlsx']:
        df = pd.read_csv(file_path) if file_type == 'csv' else pd.read_excel(file_path)
        text = df.to_string()

    # Summarize the extracted text
    summary = summarizer(text, max_length=150, min_length=50, do_sample=False)
    return summary